import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_irap_percentage_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Palette
    BG_DARK = RGBColor(11, 25, 18)       # #0B1912
    CARD_BG = RGBColor(19, 46, 35)       # #132E23
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_MUTED = RGBColor(167, 243, 208) # #A7F3D0
    ACCENT_GREEN = RGBColor(16, 185, 129)# #10B981
    ACCENT_GOLD = RGBColor(245, 158, 11) # #F59E0B
    
    img_hero = "/home/psalmprax/.gemini/antigravity-cli/brain/1c75408d-d945-4163-9d10-c51860a0efb8/smart_farm_hero_1784564335385.jpg"
    img_scanner = "/home/psalmprax/.gemini/antigravity-cli/brain/1c75408d-d945-4163-9d10-c51860a0efb8/disease_scanner_1784564346786.jpg"

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_DARK
        bg.line.color.rgb = BG_DARK
        return bg

    def add_hdr(slide, badge_txt, title_txt):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_GREEN
        badge.line.color.rgb = ACCENT_GREEN
        tf = badge.text_frame
        p = tf.paragraphs[0]
        p.text = badge_txt.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = BG_DARK
        p.alignment = PP_ALIGN.CENTER
        
        tx = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.7))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title_txt
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 1: Title
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide1)
    if os.path.exists(img_hero):
        slide1.shapes.add_picture(img_hero, Inches(6.8), Inches(0.8), Inches(5.8), Inches(5.8))
        
    tb = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "AG-EXTENSION AI"
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_GREEN
    
    p1 = tf.add_paragraph()
    p1.text = "Your AI Agronomist on Demand"
    p1.font.size = Pt(32)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "\nGiving farmers expert crop advice, instant disease scanning, and fair market prices right on their phones."
    p2.font.size = Pt(14)
    p2.font.color.rgb = TEXT_MUTED
    
    pill = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.6), Inches(5.8), Inches(0.8))
    pill.fill.solid()
    pill.fill.fore_color.rgb = CARD_BG
    pill.line.color.rgb = ACCENT_GREEN
    tf_p = pill.text_frame
    p_p = tf_p.paragraphs[0]
    p_p.text = "⚡ Seed Investment Proposal  |  NRC-IRAP R&D Aligned  |  Working MVP"
    p_p.font.size = Pt(11)
    p_p.font.bold = True
    p_p.font.color.rgb = ACCENT_GOLD

    # ----------------------------------------------------
    # SLIDE 2: Problem
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide2)
    add_hdr(slide2, "The Problem", "Farmers Face Critical Yield & Information Gaps")
    
    gaps = [
        ("1 : 3,000 Deficit", "Advisor Shortage", "Over 85% of smallholder farmers receive zero expert agronomy advice when crops fall sick."),
        ("25% - 30% Yield Loss", "Crop Destruction", "Preventable pests and plant diseases destroy up to a third of annual crop production."),
        ("30% Price Penalty", "Market Asymmetry", "Without real-time price transparency, farmers are forced to accept depressed buyer pricing.")
    ]
    for i, (val, label, desc) in enumerate(gaps):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GREEN if i == 0 else CARD_BG
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GOLD
        
        p_sub = tf.add_paragraph()
        p_sub.text = label.upper()
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT_GREEN
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 3: Solution
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide3)
    add_hdr(slide3, "The Solution", "An AI Expert in Every Farmer's Pocket")
    if os.path.exists(img_scanner):
        slide3.shapes.add_picture(img_scanner, Inches(0.8), Inches(1.8), Inches(5.0), Inches(4.8))
        
    solution_points = [
        ("OmniRoute / AIHubMix Engine:", "Score-sorted, quota-aware LLM router automatically fails over across models (Gemini 2.0, Claude 3.5, Llama 3.3) on rate limits."),
        ("2-Second Leaf Disease Diagnostics:", "Snap a photo of an unhealthy crop to receive immediate computer vision diagnosis and treatment advice."),
        ("21+ Model Context Protocol (MCP) Tools:", "Connects live telemetry from weather forecasts, satellite NDVI indices, and research databases."),
        ("Multi-Channel & Low-Bandwidth:", "Runs over WhatsApp, SMS, or voice notes — built for rural low-connectivity environments.")
    ]
    for i, (title_t, desc_t) in enumerate(solution_points):
        card = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(1.8 + i*1.3), Inches(6.4), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)
        
        p = tf.paragraphs[0]
        p.text = title_t
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc_t
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 4: Product Stage & Development Roadmap
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide4)
    add_hdr(slide4, "Product Stage", "Working Prototype Ready for Field Deployment")
    
    t_metrics = [
        ("Functional MVP Built", "Core Tech Complete", "Multi-agent engine, 21+ MCP tools, and multi-channel UI fully functional."),
        ("Pilot Partner Pipeline", "Co-op Partnerships", "In active discussions with regional agricultural co-ops for initial field deployment."),
        ("Target 90%+ Precision", "Model Calibration", "Ongoing calibration against certified plant disease datasets."),
        ("15,000 Farmer Target", "Pilot Horizon", "Capital deployment will scale field testing across 15,000 active farmers in year one.")
    ]
    cw = Inches(5.6)
    ch = Inches(2.2)
    coords = [(Inches(0.8), Inches(1.8)), (Inches(6.8), Inches(1.8)), (Inches(0.8), Inches(4.3)), (Inches(6.8), Inches(4.3))]
    
    for i, (t_val, t_lbl, t_desc) in enumerate(t_metrics):
        cx, cy = coords[i]
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GREEN
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.25)
        
        p = tf.paragraphs[0]
        p.text = t_val
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GOLD
        
        p_lbl = tf.add_paragraph()
        p_lbl.text = t_lbl.upper()
        p_lbl.font.size = Pt(11)
        p_lbl.font.bold = True
        p_lbl.font.color.rgb = ACCENT_GREEN
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n{t_desc}"
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 5: Go-To-Market Strategy
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide5)
    add_hdr(slide5, "How We Grow", "Fast & Low-Cost Customer Growth Plan")
    
    channels = [
        ("Co-op Partnerships", "Top-Down B2B", "Selling to agricultural co-ops who instantly onboard 10,000+ farmers per deployment."),
        ("Local Supply Stores", "Store Kiosks", "Agro-dealers use our diagnostic kiosks in-store to drive targeted input sales."),
        ("WhatsApp Sharing", "Word of Mouth", "Farmers share voice notes and disease reports in community groups, driving viral adoption.")
    ]
    for i, (ch_name, ch_sub, ch_desc) in enumerate(channels):
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GREEN
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = ch_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p_sub = tf.add_paragraph()
        p_sub.text = ch_sub.upper()
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT_GOLD
        
        p_desc = tf.add_paragraph()
        p_desc.text = f"\n{ch_desc}"
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 6: Unit Economics & Margin Profile
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide6)
    add_hdr(slide6, "Unit Economics", "High-Margin Software & Revenue Model")
    
    card_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = ACCENT_GREEN
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.4)
    tf_l.margin_top = Inches(0.4)
    
    p = tf_l.paragraphs[0]
    p.text = "Target Efficiency Benchmarks"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    metrics = [
        ("LTV / CAC Return Ratio", "26.6x (Target efficiency benchmark)"),
        ("Gross Profit Margin", "84% (Low AI computing cost per query)"),
        ("Net Revenue Retention", "125%+ (Projected co-op tier expansion)"),
        ("Payback Period", "< 3 Months per co-op contract")
    ]
    for mlabel, mval in metrics:
        p_m = tf_l.add_paragraph()
        p_m.text = f"\n• {mlabel}: {mval}"
        p_m.font.size = Pt(13)
        p_m.font.color.rgb = TEXT_WHITE
        
    card_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.8), Inches(5.8), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = CARD_BG
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.4)
    tf_r.margin_top = Inches(0.4)
    
    p = tf_r.paragraphs[0]
    p.text = "Revenue Stream Breakdown (%)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    streams = [
        ("1. Co-op Enterprise SaaS", "75% of Total Revenue"),
        ("2. Commercial Farm Plans", "15% of Total Revenue"),
        ("3. Data Insights API", "7% of Total Revenue"),
        ("4. Marketplace Transaction Fees", "3% of Total Revenue")
    ]
    for stitle, sdesc in streams:
        p_s = tf_r.add_paragraph()
        p_s.text = f"\n{stitle}"
        p_s.font.size = Pt(14)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_WHITE
        
        p_sd = tf_r.add_paragraph()
        p_sd.text = sdesc
        p_sd.font.size = Pt(12)
        p_sd.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 7: Financial Growth Multipliers
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide7)
    add_hdr(slide7, "Growth Trajectory", "3-Year ARR Expansion Multipliers")
    
    fin_years = [
        ("Year 1 Target", "1.0x Base Scale", "20 Co-ops | 15,000 Farmers", "Establish regional pilot footprint and validate unit economics."),
        ("Year 2 Target", "6.4x ARR Expansion", "110 Co-ops | 120,000 Farmers", "Expand across 3 countries & launch Data Insights API."),
        ("Year 3 Target", "25.0x ARR Scale", "400 Co-ops | 650,000 Farmers", "Scale nationwide and integrate with agricultural insurers.")
    ]
    for i, (y_name, y_arr, y_metrics, y_goal) in enumerate(fin_years):
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GOLD if i == 2 else ACCENT_GREEN
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        
        p = tf.paragraphs[0]
        p.text = y_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        
        p_arr = tf.add_paragraph()
        p_arr.text = y_arr
        p_arr.font.size = Pt(24)
        p_arr.font.bold = True
        p_arr.font.color.rgb = ACCENT_GOLD
        
        p_m = tf.add_paragraph()
        p_m.text = f"\n{y_metrics}"
        p_m.font.size = Pt(12)
        p_m.font.color.rgb = ACCENT_GREEN
        
        p_g = tf.add_paragraph()
        p_g.text = f"\nObjective:\n{y_goal}"
        p_g.font.size = Pt(12)
        p_g.font.color.rgb = TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 8: Market Share Potential
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide8)
    add_hdr(slide8, "Market Potential", "High-Growth AgTech Opportunity")
    
    markets = [
        ("TAM Share", "100% Market Base", "Global Smart Agriculture & Decision Support Market."),
        ("SAM Share", "27% Addressable Segment", "Digital Advisory Services in Emerging & Developing Regions."),
        ("SOM Target", "3.2% Initial Target Share", "Targetable Co-ops and Commercial Farms in initial deployment regions.")
    ]
    for i, (m_type, m_val, m_desc) in enumerate(markets):
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BG
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        
        p = tf.paragraphs[0]
        p.text = m_type
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GREEN
        
        p_v = tf.add_paragraph()
        p_v.text = m_val
        p_v.font.size = Pt(24)
        p_v.font.bold = True
        p_v.font.color.rgb = ACCENT_GOLD
        
        p_d = tf.add_paragraph()
        p_d.text = f"\n{m_desc}"
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 9: Founding Team
    # ----------------------------------------------------
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide9)
    add_hdr(slide9, "The Team", "Proven Execution & Agronomy Expertise")
    
    team = [
        ("CEO & Co-Founder", "Serial AgTech founder with 10+ years experience building and scaling sales networks."),
        ("CTO & Co-Founder", "AI Systems Architect specializing in computer vision, multi-agent AI, and smart tools."),
        ("Head of Agronomy", "Former Director of Agricultural Extension with 15+ years leading field advice teams.")
    ]
    for i, (t_role, t_bio) in enumerate(team):
        card = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*3.9), Inches(1.8), Inches(3.6), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = ACCENT_GREEN
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.4)
        
        p = tf.paragraphs[0]
        p.text = t_role
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_GOLD
        
        p_b = tf.add_paragraph()
        p_b.text = f"\n{t_bio}"
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 10: Funding Requirement & NRC-IRAP Alignment
    # ----------------------------------------------------
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(slide10)
    add_hdr(slide10, "Investment & NRC-IRAP", "Why We Need Capital & IRAP Innovation Alignment")
    
    card_l = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = CARD_BG
    card_l.line.color.rgb = ACCENT_GOLD
    tf_l = card_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.4)
    tf_l.margin_top = Inches(0.3)
    
    p = tf_l.paragraphs[0]
    p.text = "Why We Need Investment"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    needs = [
        ("1. Scale R&D & Multi-Agent AI:", "Transition working MVP to production-grade distributed edge AI."),
        ("2. Expand Co-op Field Pilots:", "Fund on-the-ground validation with 50 agricultural co-op partners."),
        ("3. High-Skill Talent Acquisition:", "Hire specialized AI systems engineers and agronomy field researchers.")
    ]
    for ntitle, ndesc in needs:
        p_n = tf_l.add_paragraph()
        p_n.text = f"\n{ntitle}"
        p_n.font.size = Pt(13)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_WHITE
        p_nd = tf_l.add_paragraph()
        p_nd.text = ndesc
        p_nd.font.size = Pt(11)
        p_nd.font.color.rgb = TEXT_MUTED

    card_r = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = CARD_BG
    card_r.line.color.rgb = ACCENT_GREEN
    tf_r = card_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.4)
    tf_r.margin_top = Inches(0.3)
    
    p = tf_r.paragraphs[0]
    p.text = "NRC-IRAP R&D Alignment"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    
    irap_points = [
        ("OCAP® Indigenous Data Sovereignty:", "Tiered consent microservices protecting First Nations data ownership while preserving AI diagnostic utility."),
        ("Canadian Regulatory Compliance Engine:", "Automated validation against federal (Fertilizers Act) and provincial regulations."),
        ("Climate Resilience Adaptation Toolkit:", "Downscaling climate models (CCCma) for Canadian agro-climatic zones (Prairie drylands, freeze-thaw cycles)."),
        ("Non-Dilutive R&D Salary Co-Funding:", "IRAP co-funds Canadian AI software engineers, data scientists & Indigenous knowledge specialists creating IP in Canada.")
    ]
    for ititle, idesc in irap_points:
        p_i = tf_r.add_paragraph()
        p_i.text = f"\n• {ititle}"
        p_i.font.size = Pt(13)
        p_i.font.bold = True
        p_i.font.color.rgb = TEXT_WHITE
        p_id = tf_r.add_paragraph()
        p_id.text = idesc
        p_id.font.size = Pt(11)
        p_id.font.color.rgb = TEXT_MUTED

    out_path = '/home/psalmprax/ALL_PROJECTS/ag_extension_decision_support/Ag_Extension_Investor_Deck.pptx'
    prs.save(out_path)
    print(f"10-Slide Percentage & IRAP Deck successfully saved at {out_path}")

if __name__ == '__main__':
    create_irap_percentage_deck()
